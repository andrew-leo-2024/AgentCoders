#!/usr/bin/env python3
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

# Colors
CHERRY = RGBColor(0x99, 0x00, 0x11)
OFF_WHITE = RGBColor(0xFC, 0xF6, 0xF5)
NAVY = RGBColor(0x2F, 0x3C, 0x7E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY = RGBColor(0x66, 0x66, 0x66)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)

def add_dark_slide():
    """Add a blank dark slide"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = NAVY
    return slide

def add_light_slide():
    """Add a blank light slide"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = OFF_WHITE
    return slide

def add_white_slide():
    """Add a blank white slide"""
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE
    return slide

def set_text_props(text_frame, font_size, font_name, bold=False, italic=False, color=WHITE):
    """Set text properties"""
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(font_size)
            run.font.name = font_name
            run.font.bold = bold
            run.font.italic = italic
            run.font.color.rgb = color

# ============================================
# SLIDE 1: TITLE
# ============================================
slide = add_dark_slide()

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "The Next 30 Days"
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = WHITE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_frame.word_wrap = True

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.9), Inches(9), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "From $0 to First Revenue — One Action at a Time"
subtitle_frame.paragraphs[0].font.size = Pt(20)
subtitle_frame.paragraphs[0].font.name = "Calibri"
subtitle_frame.paragraphs[0].font.color.rgb = OFF_WHITE
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add footer
footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.4))
footer_frame = footer_box.text_frame
footer_frame.text = "FrankMax Digital | February 2026"
footer_frame.paragraphs[0].font.size = Pt(14)
footer_frame.paragraphs[0].font.name = "Calibri"
footer_frame.paragraphs[0].font.color.rgb = WHITE
footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add accent line
line = slide.shapes.add_shape(1, Inches(3.5), Inches(1.5), Inches(3), Inches(0.08))
line.fill.solid()
line.fill.fore_color.rgb = CHERRY
line.line.color.rgb = CHERRY

# ============================================
# SLIDE 2: THE TRUTH
# ============================================
slide = add_dark_slide()
slide.background.fill.fore_color.rgb = CHERRY

# Main title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(9), Inches(0.7))
title_frame = title_box.text_frame
title_frame.text = "You have $1,000 and an idea."
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = WHITE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Truths
truth_box = slide.shapes.add_textbox(Inches(1.5), Inches(1.8), Inches(7), Inches(1.5))
truth_frame = truth_box.text_frame
truth_frame.text = "You do not have a product.\nYou do not have a customer.\nYou do not have a website."
truth_frame.paragraphs[0].font.size = Pt(20)
truth_frame.paragraphs[0].font.name = "Calibri"
truth_frame.paragraphs[0].font.color.rgb = OFF_WHITE
truth_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
for para in truth_frame.paragraphs:
    para.font.size = Pt(20)
    para.font.name = "Calibri"
    para.font.color.rgb = OFF_WHITE
    para.alignment = PP_ALIGN.CENTER

# Challenge
challenge_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.8))
challenge_frame = challenge_box.text_frame
challenge_frame.text = "You have 30 days to change that."
challenge_frame.paragraphs[0].font.size = Pt(28)
challenge_frame.paragraphs[0].font.bold = True
challenge_frame.paragraphs[0].font.name = "Georgia"
challenge_frame.paragraphs[0].font.color.rgb = WHITE
challenge_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 3: THE ONE SERVICE
# ============================================
slide = add_light_slide()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "One Service. One Buyer. One Price."
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = NAVY
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Card background
card = slide.shapes.add_shape(1, Inches(1.5), Inches(1.3), Inches(7), Inches(3.5))
card.fill.solid()
card.fill.fore_color.rgb = WHITE
card.line.color.rgb = GRAY
card.line.width = Pt(1)

# Service title
service_box = slide.shapes.add_textbox(Inches(1.7), Inches(1.5), Inches(6.6), Inches(0.5))
service_frame = service_box.text_frame
service_frame.text = "Revenue Chokepoint Diagnostic"
service_frame.paragraphs[0].font.size = Pt(28)
service_frame.paragraphs[0].font.bold = True
service_frame.paragraphs[0].font.name = "Georgia"
service_frame.paragraphs[0].font.color.rgb = CHERRY
service_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Description
desc_box = slide.shapes.add_textbox(Inches(1.9), Inches(2.1), Inches(6.2), Inches(0.5))
desc_frame = desc_box.text_frame
desc_frame.text = "Find the #1 operational bottleneck costing your company revenue."
desc_frame.paragraphs[0].font.size = Pt(16)
desc_frame.paragraphs[0].font.name = "Calibri"
desc_frame.paragraphs[0].font.color.rgb = DARK_GRAY
desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Timeline
timeline_box = slide.shapes.add_textbox(Inches(1.9), Inches(2.7), Inches(6.2), Inches(0.4))
timeline_frame = timeline_box.text_frame
timeline_frame.text = "Deliver in 5-10 business days."
timeline_frame.paragraphs[0].font.size = Pt(16)
timeline_frame.paragraphs[0].font.name = "Calibri"
timeline_frame.paragraphs[0].font.color.rgb = DARK_GRAY
timeline_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Price
price_box = slide.shapes.add_textbox(Inches(1.9), Inches(3.2), Inches(6.2), Inches(0.5))
price_frame = price_box.text_frame
price_frame.text = "Price: $5,000 - $15,000"
price_frame.paragraphs[0].font.size = Pt(24)
price_frame.paragraphs[0].font.bold = True
price_frame.paragraphs[0].font.name = "Georgia"
price_frame.paragraphs[0].font.color.rgb = NAVY
price_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Cost
cost_box = slide.shapes.add_textbox(Inches(1.9), Inches(3.8), Inches(6.2), Inches(0.4))
cost_frame = cost_box.text_frame
cost_frame.text = "Cost to deliver: $0 (your expertise + AI)"
cost_frame.paragraphs[0].font.size = Pt(14)
cost_frame.paragraphs[0].font.name = "Calibri"
cost_frame.paragraphs[0].font.color.rgb = GRAY
cost_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Margin
margin_box = slide.shapes.add_textbox(Inches(1.9), Inches(4.3), Inches(6.2), Inches(0.4))
margin_frame = margin_box.text_frame
margin_frame.text = "Margin: 100%"
margin_frame.paragraphs[0].font.size = Pt(18)
margin_frame.paragraphs[0].font.bold = True
margin_frame.paragraphs[0].font.name = "Georgia"
margin_frame.paragraphs[0].font.color.rgb = CHERRY
margin_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 4: THE BUYER
# ============================================
slide = add_white_slide()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Who Writes the Check"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = NAVY
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Three buyer cards
buyers = [
    ("COO / Founder", "$5M-$100M companies, feels pain daily, has budget authority, 7-14 day sales cycle", 0.6),
    ("CFO", "Revenue leakage, billing drift, 14-21 day cycle, $10-25K deals", 3.6),
    ("CTO / CISO", "AI liability exposure, governance gaps, 14-30 day cycle, $8-20K deals", 6.6)
]

for title, details, x_pos in buyers:
    # Card
    card = slide.shapes.add_shape(1, Inches(x_pos), Inches(1.2), Inches(2.8), Inches(3.2))
    card.fill.solid()
    card.fill.fore_color.rgb = OFF_WHITE
    card.line.color.rgb = GRAY
    card.line.width = Pt(1)

    # Title
    title_box = slide.shapes.add_textbox(Inches(x_pos + 0.15), Inches(1.4), Inches(2.5), Inches(0.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(16)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Georgia"
    title_frame.paragraphs[0].font.color.rgb = CHERRY
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Details
    details_box = slide.shapes.add_textbox(Inches(x_pos + 0.15), Inches(2.05), Inches(2.5), Inches(2.2))
    details_frame = details_box.text_frame
    details_frame.text = details
    details_frame.word_wrap = True
    details_frame.paragraphs[0].font.size = Pt(12)
    details_frame.paragraphs[0].font.name = "Calibri"
    details_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    details_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Bottom text
bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.6))
bottom_frame = bottom_box.text_frame
bottom_frame.text = "Start with COOs. They feel pain fastest and decide fastest."
bottom_frame.paragraphs[0].font.size = Pt(14)
bottom_frame.paragraphs[0].font.bold = True
bottom_frame.paragraphs[0].font.name = "Georgia"
bottom_frame.paragraphs[0].font.color.rgb = CHERRY
bottom_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 5: DAYS 1-7: EXIST
# ============================================
slide = add_light_slide()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Week 1: Exist"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = NAVY
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Seven day boxes
days1 = [
    ("D1", "Register domain ($12)"),
    ("D2", "Build one-page site (Carrd, $0)"),
    ("D3", "LinkedIn post #1"),
    ("D4", "Identify 50 target companies"),
    ("D5", "Send 20 connection requests"),
    ("D6", "Send 20 more"),
    ("D7", "Follow up, engage on posts")
]

for idx, (day, action) in enumerate(days1):
    x_pos = 0.6 + (idx * 1.3)

    # Card
    card = slide.shapes.add_shape(1, Inches(x_pos), Inches(1.3), Inches(1.15), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = GRAY
    card.line.width = Pt(1)

    # Day label
    day_box = slide.shapes.add_textbox(Inches(x_pos + 0.08), Inches(1.45), Inches(1), Inches(0.35))
    day_frame = day_box.text_frame
    day_frame.text = day
    day_frame.paragraphs[0].font.size = Pt(14)
    day_frame.paragraphs[0].font.bold = True
    day_frame.paragraphs[0].font.name = "Georgia"
    day_frame.paragraphs[0].font.color.rgb = CHERRY
    day_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Action
    action_box = slide.shapes.add_textbox(Inches(x_pos + 0.08), Inches(1.85), Inches(1), Inches(0.85))
    action_frame = action_box.text_frame
    action_frame.text = action
    action_frame.word_wrap = True
    action_frame.paragraphs[0].font.size = Pt(11)
    action_frame.paragraphs[0].font.name = "Calibri"
    action_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    action_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Stats
stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.4))
stats_frame = stats_box.text_frame
stats_frame.text = "Cost: $12 | Revenue: $0 | Conversations: 0"
stats_frame.paragraphs[0].font.size = Pt(14)
stats_frame.paragraphs[0].font.bold = True
stats_frame.paragraphs[0].font.name = "Georgia"
stats_frame.paragraphs[0].font.color.rgb = NAVY
stats_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Goal
goal_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.5))
goal_frame = goal_box.text_frame
goal_frame.text = "Goal: Be findable. Be reachable. Be credible."
goal_frame.paragraphs[0].font.size = Pt(16)
goal_frame.paragraphs[0].font.bold = True
goal_frame.paragraphs[0].font.name = "Georgia"
goal_frame.paragraphs[0].font.color.rgb = CHERRY
goal_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 6: DAYS 8-14: CONVERSATIONS
# ============================================
slide = add_white_slide()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Week 2: Conversations"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = NAVY
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Days
days2 = [
    ("D8", "LinkedIn post #2 (chokepoint example)"),
    ("D9", "Book first 3 discovery calls"),
    ("D10-11", "Conduct 3 calls, document pain"),
    ("D12", "Write 1-page diagnostic proposal"),
    ("D13", "Send 3 proposals"),
    ("D14", "Follow up, book 3 more calls")
]

for idx, (day, action) in enumerate(days2):
    x_pos = 0.6 + (idx * 1.5)
    w = 1.35

    # Card
    card = slide.shapes.add_shape(1, Inches(x_pos), Inches(1.3), Inches(w), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = OFF_WHITE
    card.line.color.rgb = GRAY
    card.line.width = Pt(1)

    # Day label
    day_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(1.45), Inches(w - 0.2), Inches(0.3))
    day_frame = day_box.text_frame
    day_frame.text = day
    day_frame.paragraphs[0].font.size = Pt(13)
    day_frame.paragraphs[0].font.bold = True
    day_frame.paragraphs[0].font.name = "Georgia"
    day_frame.paragraphs[0].font.color.rgb = CHERRY
    day_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Action
    action_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(1.8), Inches(w - 0.2), Inches(0.9))
    action_frame = action_box.text_frame
    action_frame.text = action
    action_frame.word_wrap = True
    action_frame.paragraphs[0].font.size = Pt(10)
    action_frame.paragraphs[0].font.name = "Calibri"
    action_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    action_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Stats
stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.4))
stats_frame = stats_box.text_frame
stats_frame.text = "Cost: $0 | Revenue: $0 | Conversations: 6"
stats_frame.paragraphs[0].font.size = Pt(14)
stats_frame.paragraphs[0].font.bold = True
stats_frame.paragraphs[0].font.name = "Georgia"
stats_frame.paragraphs[0].font.color.rgb = NAVY
stats_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Goal
goal_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.5))
goal_frame = goal_box.text_frame
goal_frame.text = "Goal: Hear real pain from real buyers in their own words."
goal_frame.paragraphs[0].font.size = Pt(16)
goal_frame.paragraphs[0].font.bold = True
goal_frame.paragraphs[0].font.name = "Georgia"
goal_frame.paragraphs[0].font.color.rgb = CHERRY
goal_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 7: DAYS 15-21: FIRST MONEY
# ============================================
slide = add_light_slide()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Week 3: First Money"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = NAVY
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Days
days3 = [
    ("D15-16", "Conduct calls #4-6, LinkedIn post #3", False),
    ("D17-18", "Send 3 more proposals, follow up on all 6", False),
    ("D19", "CLOSE FIRST DEAL ($5,000-$15,000)", True),
    ("D20-21", "Begin delivery, document everything", False)
]

for idx, (day, action, highlight) in enumerate(days3):
    x_pos = 1.2 + (idx * 1.8)
    w = 1.6

    # Highlight background if needed
    if highlight:
        bg = slide.shapes.add_shape(1, Inches(x_pos - 0.1), Inches(1.2), Inches(w + 0.2), Inches(1.7))
        bg.fill.solid()
        bg.fill.fore_color.rgb = CHERRY
        bg.line.color.rgb = CHERRY

    # Card
    card = slide.shapes.add_shape(1, Inches(x_pos), Inches(1.3), Inches(w), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = CHERRY if highlight else WHITE
    card.line.color.rgb = CHERRY if highlight else GRAY
    card.line.width = Pt(1)

    # Day label
    day_box = slide.shapes.add_textbox(Inches(x_pos + 0.12), Inches(1.45), Inches(w - 0.24), Inches(0.3))
    day_frame = day_box.text_frame
    day_frame.text = day
    day_frame.paragraphs[0].font.size = Pt(13)
    day_frame.paragraphs[0].font.bold = True
    day_frame.paragraphs[0].font.name = "Georgia"
    day_frame.paragraphs[0].font.color.rgb = WHITE if highlight else CHERRY
    day_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Action
    action_box = slide.shapes.add_textbox(Inches(x_pos + 0.12), Inches(1.8), Inches(w - 0.24), Inches(0.9))
    action_frame = action_box.text_frame
    action_frame.text = action
    action_frame.word_wrap = True
    action_frame.paragraphs[0].font.size = Pt(11)
    action_frame.paragraphs[0].font.name = "Calibri"
    action_frame.paragraphs[0].font.color.rgb = WHITE if highlight else DARK_GRAY
    action_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Stats
stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.4))
stats_frame = stats_box.text_frame
stats_frame.text = "Cost: $0 | Revenue: $5,000-$15,000 | Conversations: 10"
stats_frame.paragraphs[0].font.size = Pt(14)
stats_frame.paragraphs[0].font.bold = True
stats_frame.paragraphs[0].font.name = "Georgia"
stats_frame.paragraphs[0].font.color.rgb = NAVY
stats_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Goal
goal_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.5))
goal_frame = goal_box.text_frame
goal_frame.text = "Goal: Someone pays you money for solving their problem."
goal_frame.paragraphs[0].font.size = Pt(16)
goal_frame.paragraphs[0].font.bold = True
goal_frame.paragraphs[0].font.name = "Georgia"
goal_frame.paragraphs[0].font.color.rgb = CHERRY
goal_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 8: DAYS 22-30: COMPOUND
# ============================================
slide = add_white_slide()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "Week 4: Compound"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = NAVY
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Days
days4 = [
    ("D22", "Present findings to first client"),
    ("D23", "Close second deal"),
    ("D24-25", "Deliver + publish case study (anonymized)"),
    ("D26-27", "10 new prospects from credibility"),
    ("D28", "Propose retainer to first client ($2-5K/mo)"),
    ("D29-30", "Revenue review")
]

for idx, (day, action) in enumerate(days4):
    x_pos = 0.6 + (idx * 1.5)
    w = 1.35

    # Card
    card = slide.shapes.add_shape(1, Inches(x_pos), Inches(1.3), Inches(w), Inches(1.5))
    card.fill.solid()
    card.fill.fore_color.rgb = OFF_WHITE
    card.line.color.rgb = GRAY
    card.line.width = Pt(1)

    # Day label
    day_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(1.45), Inches(w - 0.2), Inches(0.3))
    day_frame = day_box.text_frame
    day_frame.text = day
    day_frame.paragraphs[0].font.size = Pt(13)
    day_frame.paragraphs[0].font.bold = True
    day_frame.paragraphs[0].font.name = "Georgia"
    day_frame.paragraphs[0].font.color.rgb = CHERRY
    day_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Action
    action_box = slide.shapes.add_textbox(Inches(x_pos + 0.1), Inches(1.8), Inches(w - 0.2), Inches(0.9))
    action_frame = action_box.text_frame
    action_frame.text = action
    action_frame.word_wrap = True
    action_frame.paragraphs[0].font.size = Pt(10)
    action_frame.paragraphs[0].font.name = "Calibri"
    action_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    action_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Stats
stats_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.1), Inches(9), Inches(0.4))
stats_frame = stats_box.text_frame
stats_frame.text = "Cost: $0 | Revenue: $10,000-$30,000 | Conversations: 15+"
stats_frame.paragraphs[0].font.size = Pt(14)
stats_frame.paragraphs[0].font.bold = True
stats_frame.paragraphs[0].font.name = "Georgia"
stats_frame.paragraphs[0].font.color.rgb = NAVY
stats_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Goal
goal_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.6), Inches(9), Inches(0.5))
goal_frame = goal_box.text_frame
goal_frame.text = "Goal: Prove the machine works. Then run it again."
goal_frame.paragraphs[0].font.size = Pt(16)
goal_frame.paragraphs[0].font.bold = True
goal_frame.paragraphs[0].font.name = "Georgia"
goal_frame.paragraphs[0].font.color.rgb = CHERRY
goal_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 9: THE MATH
# ============================================
slide = add_light_slide()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "30-Day Financial Reality"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = NAVY
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Left column header
left_header = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(4.2), Inches(0.5))
left_header_frame = left_header.text_frame
left_header_frame.text = "Investment"
left_header_frame.paragraphs[0].font.size = Pt(18)
left_header_frame.paragraphs[0].font.bold = True
left_header_frame.paragraphs[0].font.name = "Georgia"
left_header_frame.paragraphs[0].font.color.rgb = CHERRY
left_header_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Left column items
left_items = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(4), Inches(1.5))
left_items_frame = left_items.text_frame
left_items_frame.text = "Total capital deployed: <$50\nTools: $0-$20/mo\nTime: 30 days full-time"
left_items_frame.word_wrap = True
left_items_frame.paragraphs[0].font.size = Pt(14)
left_items_frame.paragraphs[0].font.name = "Calibri"
left_items_frame.paragraphs[0].font.color.rgb = DARK_GRAY
for para in left_items_frame.paragraphs:
    para.font.size = Pt(14)
    para.font.name = "Calibri"
    para.font.color.rgb = DARK_GRAY
    para.alignment = PP_ALIGN.LEFT

# Divider
divider = slide.shapes.add_shape(1, Inches(4.8), Inches(1.2), Inches(0.08), Inches(3.2))
divider.fill.solid()
divider.fill.fore_color.rgb = NAVY
divider.line.color.rgb = NAVY

# Right column header
right_header = slide.shapes.add_textbox(Inches(5.2), Inches(1.2), Inches(4.2), Inches(0.5))
right_header_frame = right_header.text_frame
right_header_frame.text = "Return"
right_header_frame.paragraphs[0].font.size = Pt(18)
right_header_frame.paragraphs[0].font.bold = True
right_header_frame.paragraphs[0].font.name = "Georgia"
right_header_frame.paragraphs[0].font.color.rgb = CHERRY
right_header_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Right column items
right_items = slide.shapes.add_textbox(Inches(5.4), Inches(1.8), Inches(4), Inches(2))
right_items_frame = right_items.text_frame
right_items_frame.text = "Deals closed: 2-3\nRevenue: $10,000-$30,000\nPipeline: 5-10 proposals\nConversations: 15+\nCase studies: 1-2"
right_items_frame.word_wrap = True
right_items_frame.paragraphs[0].font.size = Pt(14)
right_items_frame.paragraphs[0].font.name = "Calibri"
right_items_frame.paragraphs[0].font.color.rgb = DARK_GRAY
for para in right_items_frame.paragraphs:
    para.font.size = Pt(14)
    para.font.name = "Calibri"
    para.font.color.rgb = DARK_GRAY
    para.alignment = PP_ALIGN.LEFT

# Bottom ROI
roi_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.6))
roi_frame = roi_box.text_frame
roi_frame.text = "ROI: 200x-600x on capital. The only investment is your time."
roi_frame.paragraphs[0].font.size = Pt(14)
roi_frame.paragraphs[0].font.bold = True
roi_frame.paragraphs[0].font.name = "Georgia"
roi_frame.paragraphs[0].font.color.rgb = CHERRY
roi_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 10: THE UPSELL LADDER
# ============================================
slide = add_white_slide()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "How $5K Becomes $500K"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = NAVY
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Upsell steps
upsells = [
    ("Step 1", "Diagnostic Sprint", "$5-15K (one-time)"),
    ("Step 2", "Fix Implementation", "$15-30K (project)"),
    ("Step 3", "Monthly Monitoring Retainer", "$2-5K/mo ($24-60K/yr)"),
    ("Step 4", "Governance Assessment", "$8-20K"),
    ("Step 5", "Frankmax PIAR", "$25-75K (institutional)"),
    ("Step 6", "Enterprise License", "$100-500K/yr")
]

for idx, (step, title, price) in enumerate(upsells):
    y_start = 1.3 + (idx * 0.55)
    step_width = 0.8 + (idx * 0.12)

    # Card
    card = slide.shapes.add_shape(1, Inches(1.0 + (idx * 0.15)), Inches(y_start), Inches(step_width), Inches(0.5))
    card.fill.solid()
    card.fill.fore_color.rgb = NAVY
    card.line.color.rgb = NAVY

    # Title
    title_box = slide.shapes.add_textbox(Inches(1.8 + (idx * 0.15)), Inches(y_start + 0.05), Inches(5), Inches(0.4))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(13)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Georgia"
    title_frame.paragraphs[0].font.color.rgb = CHERRY
    title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Price
    price_box = slide.shapes.add_textbox(Inches(1.8 + (idx * 0.15)), Inches(y_start + 0.25), Inches(5), Inches(0.2))
    price_frame = price_box.text_frame
    price_frame.text = price
    price_frame.paragraphs[0].font.size = Pt(11)
    price_frame.paragraphs[0].font.name = "Calibri"
    price_frame.paragraphs[0].font.color.rgb = GRAY
    price_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Bottom text
bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.4))
bottom_frame = bottom_box.text_frame
bottom_frame.text = "One client can become $200K+ over 18 months"
bottom_frame.paragraphs[0].font.size = Pt(15)
bottom_frame.paragraphs[0].font.bold = True
bottom_frame.paragraphs[0].font.name = "Georgia"
bottom_frame.paragraphs[0].font.color.rgb = CHERRY
bottom_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 11: WHAT NOT TO DO
# ============================================
slide = add_dark_slide()
slide.background.fill.fore_color.rgb = NAVY

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "The 10 Traps"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = WHITE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Traps
traps = [
    "Do NOT build a product (sell the service first)",
    "Do NOT write more strategy documents",
    "Do NOT design a logo or brand identity",
    "Do NOT build the 74-system architecture",
    "Do NOT incorporate in multiple jurisdictions",
    "Do NOT attend conferences",
    "Do NOT hire anyone",
    "Do NOT raise capital",
    "Do NOT build an app",
    "Do NOT think about the ORF protocol until $50K revenue"
]

left_traps = traps[:5]
right_traps = traps[5:]

for idx, trap in enumerate(left_traps):
    y_pos = 1.2 + (idx * 0.7)

    trap_box = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos), Inches(4.3), Inches(0.65))
    trap_frame = trap_box.text_frame
    trap_frame.word_wrap = True

    # Add number and text
    trap_frame.text = f"{idx + 1}. {trap[4:]}"
    trap_frame.paragraphs[0].font.size = Pt(13)
    trap_frame.paragraphs[0].font.name = "Calibri"
    trap_frame.paragraphs[0].font.color.rgb = OFF_WHITE
    trap_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

for idx, trap in enumerate(right_traps):
    y_pos = 1.2 + (idx * 0.7)

    trap_box = slide.shapes.add_textbox(Inches(5.3), Inches(y_pos), Inches(4.3), Inches(0.65))
    trap_frame = trap_box.text_frame
    trap_frame.word_wrap = True

    # Add number and text
    trap_frame.text = f"{idx + 6}. {trap[4:]}"
    trap_frame.paragraphs[0].font.size = Pt(13)
    trap_frame.paragraphs[0].font.name = "Calibri"
    trap_frame.paragraphs[0].font.color.rgb = OFF_WHITE
    trap_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# ============================================
# SLIDE 12: THE CONVERSATION SCRIPT
# ============================================
slide = add_light_slide()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "The Only Script You Need"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = NAVY
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Script steps
scripts = [
    ("1", "CONNECT", "I help [industry] companies find their biggest revenue bottleneck."),
    ("2", "ASK", "What's the single biggest operational issue costing you money right now?"),
    ("3", "QUANTIFY", "If we fixed that, what would it be worth in monthly revenue?"),
    ("4", "OFFER", "I can identify your top 3 chokepoints with dollar values in 5-10 days. $[X]. If I don't find 3x that amount, I refund you.")
]

for idx, (num, label, text) in enumerate(scripts):
    y_pos = 1.2 + (idx * 0.9)

    # Circle
    circle = slide.shapes.add_shape(1, Inches(0.8), Inches(y_pos + 0.05), Inches(0.4), Inches(0.4))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CHERRY
    circle.line.color.rgb = CHERRY

    # Number
    num_box = slide.shapes.add_textbox(Inches(0.8), Inches(y_pos + 0.05), Inches(0.4), Inches(0.4))
    num_frame = num_box.text_frame
    num_frame.text = num
    num_frame.paragraphs[0].font.size = Pt(18)
    num_frame.paragraphs[0].font.bold = True
    num_frame.paragraphs[0].font.name = "Georgia"
    num_frame.paragraphs[0].font.color.rgb = WHITE
    num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    num_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Label
    label_box = slide.shapes.add_textbox(Inches(1.4), Inches(y_pos + 0.08), Inches(1.2), Inches(0.3))
    label_frame = label_box.text_frame
    label_frame.text = label
    label_frame.paragraphs[0].font.size = Pt(13)
    label_frame.paragraphs[0].font.bold = True
    label_frame.paragraphs[0].font.name = "Georgia"
    label_frame.paragraphs[0].font.color.rgb = CHERRY
    label_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

    # Text
    text_box = slide.shapes.add_textbox(Inches(1.4), Inches(y_pos + 0.4), Inches(7.6), Inches(0.5))
    text_frame = text_box.text_frame
    text_frame.text = text
    text_frame.word_wrap = True
    text_frame.paragraphs[0].font.size = Pt(12)
    text_frame.paragraphs[0].font.name = "Calibri"
    text_frame.paragraphs[0].font.color.rgb = DARK_GRAY
    text_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# ============================================
# SLIDE 13: 90-DAY HORIZON
# ============================================
slide = add_white_slide()

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "If Day 30 Works, Then..."
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = NAVY
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Three columns
horizon = [
    ("Month 2", ["5 clients", "$25-50K", "hire first operator", "systemize delivery"]),
    ("Month 3", ["10 clients", "$50-100K", "launch Operator Track", "case study library"]),
    ("Month 4-6", ["20 clients", "$100-200K", "vertical specialization", "first retainer portfolio"])
]

for idx, (month, items) in enumerate(horizon):
    x_pos = 0.8 + (idx * 2.9)

    # Card
    card = slide.shapes.add_shape(1, Inches(x_pos), Inches(1.2), Inches(2.6), Inches(3.2))
    card.fill.solid()
    card.fill.fore_color.rgb = OFF_WHITE
    card.line.color.rgb = GRAY
    card.line.width = Pt(1)

    # Month header
    month_box = slide.shapes.add_textbox(Inches(x_pos + 0.15), Inches(1.35), Inches(2.3), Inches(0.4))
    month_frame = month_box.text_frame
    month_frame.text = month
    month_frame.paragraphs[0].font.size = Pt(15)
    month_frame.paragraphs[0].font.bold = True
    month_frame.paragraphs[0].font.name = "Georgia"
    month_frame.paragraphs[0].font.color.rgb = CHERRY
    month_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Items
    for item_idx, item in enumerate(items):
        item_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(1.9 + (item_idx * 0.6)), Inches(2.2), Inches(0.5))
        item_frame = item_box.text_frame
        item_frame.text = item
        item_frame.word_wrap = True
        item_frame.paragraphs[0].font.size = Pt(12)
        item_frame.paragraphs[0].font.name = "Calibri"
        item_frame.paragraphs[0].font.color.rgb = DARK_GRAY
        item_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Bottom text
bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.7), Inches(9), Inches(0.5))
bottom_frame = bottom_box.text_frame
bottom_frame.text = "Each month proves the next. Nothing is assumed."
bottom_frame.paragraphs[0].font.size = Pt(14)
bottom_frame.paragraphs[0].font.bold = True
bottom_frame.paragraphs[0].font.name = "Georgia"
bottom_frame.paragraphs[0].font.color.rgb = CHERRY
bottom_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 14: THE TRILLION PATH
# ============================================
slide = add_dark_slide()
slide.background.fill.fore_color.rgb = NAVY

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.6))
title_frame = title_box.text_frame
title_frame.text = "The Honest Compounding Path"
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.name = "Georgia"
title_frame.paragraphs[0].font.color.rgb = WHITE
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Milestones
milestones = [
    ("2026", "$250K-$1.2M", "prove the model, build credibility"),
    ("2027", "$3-10M", "verticalize, operator army, first institutional clients"),
    ("2028", "$20-50M", "ORF protocol deployment, cross-border, insurance layer"),
    ("2029", "$100-500M", "institutional infrastructure, government pilots"),
    ("2030", "$1-5B", "obligation infrastructure at scale, GDP-adjacent positioning")
]

for idx, (year, value, desc) in enumerate(milestones):
    y_pos = 1.3 + (idx * 0.75)

    # Year box
    year_box = slide.shapes.add_shape(1, Inches(0.7), Inches(y_pos), Inches(1.0), Inches(0.6))
    year_box.fill.solid()
    year_box.fill.fore_color.rgb = CHERRY
    year_box.line.color.rgb = CHERRY

    # Year text
    year_text = slide.shapes.add_textbox(Inches(0.7), Inches(y_pos + 0.12), Inches(1.0), Inches(0.35))
    year_frame = year_text.text_frame
    year_frame.text = year
    year_frame.paragraphs[0].font.size = Pt(13)
    year_frame.paragraphs[0].font.bold = True
    year_frame.paragraphs[0].font.name = "Georgia"
    year_frame.paragraphs[0].font.color.rgb = WHITE
    year_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    year_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Value
    value_box = slide.shapes.add_textbox(Inches(1.85), Inches(y_pos + 0.08), Inches(2.2), Inches(0.35))
    value_frame = value_box.text_frame
    value_frame.text = value
    value_frame.paragraphs[0].font.size = Pt(13)
    value_frame.paragraphs[0].font.bold = True
    value_frame.paragraphs[0].font.name = "Georgia"
    value_frame.paragraphs[0].font.color.rgb = CHERRY
    value_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    value_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Description
    desc_box = slide.shapes.add_textbox(Inches(4.2), Inches(y_pos + 0.08), Inches(5.3), Inches(0.44))
    desc_frame = desc_box.text_frame
    desc_frame.text = desc
    desc_frame.word_wrap = True
    desc_frame.paragraphs[0].font.size = Pt(12)
    desc_frame.paragraphs[0].font.name = "Calibri"
    desc_frame.paragraphs[0].font.color.rgb = OFF_WHITE
    desc_frame.paragraphs[0].alignment = PP_ALIGN.LEFT
    desc_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

# Bottom text
bottom_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.0), Inches(9), Inches(0.5))
bottom_frame = bottom_box.text_frame
bottom_frame.text = "Trillions require decades. Billions require years. Millions require months. The first $50K requires 30 days."
bottom_frame.paragraphs[0].font.size = Pt(14)
bottom_frame.paragraphs[0].font.italic = True
bottom_frame.paragraphs[0].font.name = "Calibri"
bottom_frame.paragraphs[0].font.color.rgb = OFF_WHITE
bottom_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 15: THE NEXT SINGLE ACTION
# ============================================
slide = add_dark_slide()
slide.background.fill.fore_color.rgb = CHERRY

# Main action
action_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.8))
action_frame = action_box.text_frame
action_frame.text = "Open LinkedIn."
action_frame.paragraphs[0].font.size = Pt(48)
action_frame.paragraphs[0].font.bold = True
action_frame.paragraphs[0].font.name = "Georgia"
action_frame.paragraphs[0].font.color.rgb = WHITE
action_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Next steps
steps_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(9), Inches(2))
steps_frame = steps_box.text_frame
steps_frame.text = "Find one COO at a $10M company.\nSend one message.\nThat is the next action."
steps_frame.word_wrap = True
steps_frame.paragraphs[0].font.size = Pt(28)
steps_frame.paragraphs[0].font.name = "Calibri"
steps_frame.paragraphs[0].font.color.rgb = OFF_WHITE
steps_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
for para in steps_frame.paragraphs:
    para.font.size = Pt(28)
    para.font.name = "Calibri"
    para.font.color.rgb = OFF_WHITE
    para.alignment = PP_ALIGN.CENTER

# Bottom motto
motto_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.6), Inches(9), Inches(0.6))
motto_frame = motto_box.text_frame
motto_frame.text = "Everything else is procrastination."
motto_frame.paragraphs[0].font.size = Pt(16)
motto_frame.paragraphs[0].font.italic = True
motto_frame.paragraphs[0].font.name = "Calibri"
motto_frame.paragraphs[0].font.color.rgb = OFF_WHITE
motto_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# Write file
# ============================================
prs.save("/sessions/trusting-modest-mccarthy/mnt/Invention Tools/NEXT_30_DAYS.pptx")
print("Deck created successfully!")
